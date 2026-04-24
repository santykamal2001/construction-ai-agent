"""
Multi-format text extractor.
Handles: TXT, PDF, DOCX, XLSX, CSV, PPTX, RTF, MD
"""
import csv
import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_text(file_path: str) -> str:
    """
    Extract text from a file based on its extension.
    Returns extracted text as a string, or empty string on failure.
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    try:
        if ext in (".txt", ".md", ".rtf"):
            return _extract_text_file(path)
        elif ext == ".pdf":
            return _extract_pdf(path)
        elif ext in (".docx", ".doc"):
            return _extract_docx(path)
        elif ext in (".xlsx", ".xls"):
            return _extract_excel(path)
        elif ext == ".csv":
            return _extract_csv(path)
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(path)
        else:
            logger.warning(f"Unsupported file type: {ext} — {path.name}")
            return ""
    except Exception as e:
        logger.error(f"Failed to extract {path.name}: {e}")
        return ""


def _extract_text_file(path: Path) -> str:
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return ""


def _extract_pdf(path: Path) -> str:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(path))
        pages = []
        for page_num, page in enumerate(doc, 1):
            text = page.get_text().strip()
            if text:
                pages.append(f"[Page {page_num}]\n{text}")
            else:
                # Scanned/image page — run OCR
                ocr_text = _ocr_page(page, page_num)
                if ocr_text:
                    pages.append(f"[Page {page_num} - OCR]\n{ocr_text}")
        doc.close()
        return "\n\n".join(pages)
    except ImportError:
        # Fallback to pypdf (no OCR support in this path)
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            pages = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"[Page {i}]\n{text}")
            return "\n\n".join(pages)
        except ImportError:
            logger.error("No PDF library found. Install pymupdf or pypdf.")
            return ""


def _ocr_page(page, page_num: int) -> str:
    """
    Run PaddleOCR on a single PyMuPDF page.
    Renders page to a numpy image array, then extracts text.
    """
    try:
        import numpy as np
        import fitz

        # Render page at 1.5x zoom — good OCR quality, manageable image size for CPU
        pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))

        # Convert to numpy array (RGB)
        img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
            pix.height, pix.width, pix.n
        )
        if pix.n == 4:  # RGBA → RGB
            img_array = img_array[:, :, :3]

        # Initialize PaddleOCR (downloads models on first call, cached after)
        ocr = _get_paddle_ocr()
        # Disable heavy preprocessing (doc unwarping is GPU-intended, very slow on CPU)
        results = ocr.predict(
            img_array,
            use_doc_orientation_classify=False,
            use_doc_unwarping=False
        )

        if not results:
            return ""

        lines = []
        for res in results:
            texts = res.get("rec_texts", [])
            scores = res.get("rec_scores", [])
            for text, score in zip(texts, scores):
                if score >= 0.6 and text.strip():  # filter low-confidence noise
                    lines.append(text.strip())

        return "\n".join(lines)

    except ImportError:
        logger.warning("PaddleOCR not installed. Skipping OCR for scanned pages. Run: pip install paddlepaddle paddleocr")
        return ""
    except Exception as e:
        logger.error(f"OCR failed on page {page_num}: {e}")
        return ""


_paddle_ocr_instance = None

def _get_paddle_ocr():
    """Singleton — initializes PaddleOCR once and reuses across pages."""
    global _paddle_ocr_instance
    if _paddle_ocr_instance is None:
        from paddleocr import PaddleOCR
        # Use mobile models — 10x faster on CPU vs server models, good enough for printed inspection docs
        _paddle_ocr_instance = PaddleOCR(
            text_detection_model_name="PP-OCRv4_mobile_det",
            text_recognition_model_name="PP-OCRv4_mobile_rec",
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
        )
    return _paddle_ocr_instance


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except ImportError:
        logger.error("python-docx not installed.")
        return ""


def _extract_excel(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip(" |"):
                    rows.append(row_text)
            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)
    except ImportError:
        logger.error("openpyxl not installed.")
        return ""


def _extract_csv(path: Path) -> str:
    rows = []
    for encoding in ("utf-8", "latin-1"):
        try:
            with open(path, encoding=encoding, newline="") as f:
                reader = csv.reader(f)
                for row in reader:
                    rows.append(" | ".join(row))
            return "\n".join(rows)
        except UnicodeDecodeError:
            continue
    return ""


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        logger.error("python-pptx not installed.")
        return ""
